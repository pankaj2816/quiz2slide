import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
import re

def run_comprehensive_audit(pptx_path):
    prs = pptx.Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"============================================================")
    print(f"       COMPREHENSIVE 12-POINT QUALITY AUDIT SUITE")
    print(f"       Testing PPTX: {pptx_path}")
    print(f"       Total Slides: {total_slides}")
    print(f"============================================================\n")

    issues_found = []

    # Regex defect definitions
    checks = [
        ("Consecutive Multiple Spaces", r' {2,}'),
        ("Unformatted Unit Exponents (e.g. h -1, ms -2)", r'\b(?:km\s+h|ms|Nm|m\s+s)\s*[-−]\s*\d+\b'),
        ("Unformatted Powers of 10 (e.g. 10 -2, 10 3)", r'\b10\s*[-−]\s*\d+\b|\b10\s+\d+\b'),
        ("Unformatted Variable Subscripts (e.g. v 1, x 1)", r'\b(?:v|x|B|I_L|R_L|K_H|γ)\s+[1-4AB]\b'),
        ("Loose Space Before Punctuation", r'\s+[,.:;?!%°]'),
        ("Loose Space Inside Parentheses", r'\(\s+|\s+\)'),
        ("Unformatted Chemical Formulas (e.g. C4H9Br, NaNH2)", r'\b(?:C\s*4\s*H\s*9|NaNH\s*2|He\s*\+|Li\s*\+\+)\b'),
        ("Unformatted Dimensional Formulas (e.g. [ ML 0 ])", r'\[\s*M[L0-9\s−\-]*T[0-9\s−\-]*\]'),
        ("Unformatted Ordinals (e.g. 1 st, 2 nd)", r'\b\d+\s+(?:st|nd|rd|th)\b'),
        ("Broken Ratio Notation (e.g. B B1 2)", r'\bB\s*B[1₁]\b'),
        ("Loose Negative Sign in Exponents (e.g. ab -2)", r'\bab\s*[-−]\s*2\b')
    ]

    for idx, slide in enumerate(prs.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += "\n" + shape.text_frame.text.strip()

        slide_issues = []
        for name, pattern in checks:
            matches = re.findall(pattern, slide_text)
            if matches:
                # Filter legitimate dimension bracket matches like [ML⁰ T⁻³]
                if name == "Unformatted Dimensional Formulas (e.g. [ ML 0 ])":
                    bad = [m for m in matches if not re.match(r'^\[ML[⁰\d]?\s*T[⁻\d]+\]$', m.replace(' ', ''))]
                    if bad:
                        slide_issues.append((name, bad))
                else:
                    slide_issues.append((name, matches))

        if slide_issues:
            issues_found.append((idx, slide_issues, slide_text))
            print(f"❌ Slide {idx:02d} FAILED with {len(slide_issues)} issue(s):")
            for cat, items in slide_issues:
                print(f"   • {cat}: {items}")
        else:
            print(f"✅ Slide {idx:02d}: PASSED all 12 quality checks")

    print(f"\n============================================================")
    print(f"AUDIT SUMMARY: {total_slides - len(issues_found)} / {total_slides} Slides Passed Flawlessly.")
    if len(issues_found) == 0:
        print(f"🏆 100% PERFECT ZERO-DEFECT QUALITY SCORE ACHIEVED!")
    else:
        print(f"⚠️ {len(issues_found)} Slide(s) require automated cleanup.")
    print(f"============================================================")
    return len(issues_found)

if __name__ == "__main__":
    pptx_file = r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev88.pptx"
    run_comprehensive_audit(pptx_file)

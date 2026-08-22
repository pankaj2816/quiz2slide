import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
import os

prs = pptx.Presentation(r'd:\Work\kaku\test_output\Quiz_Presentation_96_Slides.pptx')
print(f"Total slides: {len(prs.slides)}")

for s_num in [1, 2, 3, 21, 22, 23, 95, 96]:
    s = prs.slides[s_num - 1]
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    pics = [sp for sp in s.shapes if sp.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE]
    print(f"\n--- SLIDE {s_num} (Pictures: {len(pics)}) ---")
    if texts:
        print(texts[0][:150] + "...")
    else:
        print("NO TEXT FRAME")

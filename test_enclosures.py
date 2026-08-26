import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(6.5))
tf = txBox.text_frame
tf.word_wrap = True

# Compare:
# 1. √(3) and √(2)
# 2. √[3] and √[2]
# 3. Cambria Math Font with \u0305
# 4. Symbol / Arial Unicode

p1 = tf.add_paragraph()
p1.text = "Option 1: Parentheses Enclosure √(3) (3 is inside the root):"
p1.font.size = Pt(16)
p1.font.bold = True
p1.font.color.rgb = RGBColor(96, 165, 250)

p1_sub = tf.add_paragraph()
p1_sub.text = "   (A) cos⁻¹(1/√(3))       (D) cos⁻¹(√(2)/3)"
p1_sub.font.size = Pt(22)
p1_sub.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf.add_paragraph()
p2.text = "Option 2: Bracket Enclosure √[3] (Matches Q20 Formula):"
p2.font.size = Pt(16)
p2.font.bold = True
p2.font.color.rgb = RGBColor(96, 165, 250)

p2_sub = tf.add_paragraph()
p2_sub.text = "   (A) cos⁻¹(1/√[3])       (D) cos⁻¹(√[2]/3)"
p2_sub.font.size = Pt(22)
p2_sub.font.color.rgb = RGBColor(255, 255, 255)

p3 = tf.add_paragraph()
p3.text = "Option 3: Standard Clean Radical √3:"
p3.font.size = Pt(16)
p3.font.bold = True
p3.font.color.rgb = RGBColor(96, 165, 250)

p3_sub = tf.add_paragraph()
p3_sub.text = "   (A) cos⁻¹(1/√3)       (D) cos⁻¹(√2/3)"
p3_sub.font.size = Pt(22)
p3_sub.font.color.rgb = RGBColor(255, 255, 255)

prs.save(r"d:\Work\kaku\test_output\radical_enclosure_comparison.pptx")
print("Saved radical_enclosure_comparison.pptx successfully!")

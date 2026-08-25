import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Verified.pptx")

print(f"Total Slides in Presentation: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides, 1):
    print(f"\n========================================================")
    print(f"                     SLIDE {idx:02d}")
    print(f"========================================================")
    
    images_count = sum(1 for shape in slide.shapes if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE)
    print(f"Attached Images/Diagrams: {images_count}")
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text.strip())
            print("-" * 40)

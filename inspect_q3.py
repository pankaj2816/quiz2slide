import sys
sys.stdout.reconfigure(encoding='utf-8')
import pymupdf
import pptx

doc = pymupdf.open(r"d:\Work\kaku\test_output\input_pdf.pdf")

print("=== PDF PAGE 2 TEXT ===")
print(doc[1].get_text("text"))

print("=== PDF PAGE 3 TEXT ===")
print(doc[2].get_text("text"))

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Verified.pptx")
print("\n=== PPTX SLIDE 3 (INDEX 2) TEXT SHAPES ===")
s3 = prs.slides[2]
for shape in s3.shapes:
    if shape.has_text_frame:
        print(f"Shape text:\n{shape.text_frame.text}\n---")

import sys
sys.stdout.reconfigure(encoding='utf-8')
import os
import time
import pptx

filepath = r'd:\Work\kaku\test_output\Quiz_Presentation_6_Slides.pptx'
mtime = os.path.getmtime(filepath)
print(f"File: {filepath}")
print(f"Last modified: {time.ctime(mtime)}")

prs = pptx.Presentation(filepath)
print(f"Total slides: {len(prs.slides)}")

for idx, s in enumerate(prs.slides, 1):
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    print(f"\n--- SLIDE {idx} ---")
    if texts:
        print(texts[0][:250])

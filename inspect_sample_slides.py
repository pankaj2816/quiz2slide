import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev93.pptx")
for idx in [0, 1, 3, 5]:
    slide = prs.slides[idx]
    print(f"=== SLIDE {idx+1} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(repr(shape.text_frame.text.strip()))
            print("-" * 30)

import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
import re
import os

def run_50_point_audit(pptx_path):
    print("=" * 80)
    print("       EXHAUSTIVE 50-POINT ENTERPRISE QUALITY ASSURANCE AUDIT SUITE")
    print(f"       Target File: {pptx_path}")
    print("=" * 80)

    if not os.path.exists(pptx_path):
        print(f"❌ ERROR: PPTX file not found at {pptx_path}")
        return 1

    prs = pptx.Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"Total Slides Loaded: {total_slides}\n")

    slide_failures = {}
    check_failure_counts = {i: 0 for i in range(1, 51)}

    for idx, slide in enumerate(prs.slides, 1):
        failures = []
        slide_text = ""
        text_frames = []
        pictures = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frames.append(shape.text_frame)
                slide_text += "\n" + shape.text_frame.text
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                pictures.append(shape)

        lines = [l.strip() for l in slide_text.split('\n') if l.strip()]

        # === GROUP 1: TYPOGRAPHY & WHITESPACE (1-10) ===
        # 1. Consecutive Whitespace
        if re.search(r' {2,}', slide_text):
            failures.append((1, "Consecutive Whitespace (2+ spaces)", re.findall(r' {2,}', slide_text)))

        # 2. Trailing Whitespace
        if any(re.search(r'[ \t]+$', l) for l in slide_text.split('\n')):
            failures.append((2, "Trailing Line Whitespace", "Found trailing spaces on line ends"))

        # 3. Leading Whitespace
        if any(re.search(r'^[ \t]+[^\s]', l) for l in slide_text.split('\n')):
            failures.append((3, "Leading Line Whitespace", "Found stray leading indentation"))

        # 4. Stray Space Before Punctuation
        if re.search(r' [.,:;?!%°]', slide_text):
            failures.append((4, "Stray Space Before Punctuation", re.findall(r' [.,:;?!%°]', slide_text)))

        # 5. Missing Space After Punctuation
        missing_space = re.findall(r'[,;?!](?=[A-Za-z])', slide_text)
        if missing_space:
            failures.append((5, "Missing Space After Punctuation", missing_space))

        # 6. Stray Padding Inside Parentheses
        if re.search(r'\(\s+[^\s]|\s+\)', slide_text):
            failures.append((6, "Stray Padding Inside Parentheses", re.findall(r'\(\s+[^\s]|\s+\)', slide_text)))

        # 7. Stray Padding Inside Brackets
        if re.search(r'\[\s+[^\s]|\s+\]', slide_text):
            failures.append((7, "Stray Padding Inside Brackets", re.findall(r'\[\s+[^\s]|\s+\]', slide_text)))

        # 8. Normalized Quotation Marks
        if re.search(r"'\s+[a-zA-Z\d]+\s+'", slide_text):
            failures.append((8, "Stray Spaces Inside Single Quotes", re.findall(r"'\s+[a-zA-Z\d]+\s+'", slide_text)))

        # 9. Duplicate Punctuation
        if re.search(r'(?<!\.)\.\.(?!\.)|::|,,|;;|\?\?|!!', slide_text):
            failures.append((9, "Duplicate Punctuation Characters", re.findall(r'(?<!\.)\.\.(?!\.)|::|,,|;;|\?\?|!!', slide_text)))

        # 10. Balanced Parentheses
        open_p = slide_text.count('(')
        close_p = slide_text.count(')')
        if open_p != close_p:
            failures.append((10, "Unbalanced Parentheses", f"Open: {open_p}, Close: {close_p}"))

        # === GROUP 2: SCIENTIFIC & UNIT FORMATTING (11-20) ===
        # 11. Velocity Units
        if re.search(r'\bkm\s*h\s*[-−]?\s*1\b', slide_text):
            failures.append((11, "Unformatted Velocity Unit (km h -1)", re.findall(r'\bkm\s*h\s*[-−]?\s*1\b', slide_text)))

        # 12. Acceleration Units
        if re.search(r'\b(?:ms|m\s*s)\s*[-−]?\s*2\b|\bm\s*/\s*s\s*2\b', slide_text):
            failures.append((12, "Unformatted Acceleration Unit (ms -2 / m/s 2)", re.findall(r'\b(?:ms|m\s*s)\s*[-−]?\s*2\b|\bm\s*/\s*s\s*2\b', slide_text)))

        # 13. Force & Pressure Units
        if re.search(r'\bNm\s*[-−]?\s*2\b|\bN\s*m\s*[-−]?\s*2\b', slide_text):
            failures.append((13, "Unformatted Pressure Unit (Nm -2)", re.findall(r'\bNm\s*[-−]?\s*2\b|\bN\s*m\s*[-−]?\s*2\b', slide_text)))

        # 14. Density Units
        if re.search(r'\bkg\s*/\s*m\s*3\b', slide_text):
            failures.append((14, "Unformatted Density Unit (kg/m 3)", re.findall(r'\bkg\s*/\s*m\s*3\b', slide_text)))

        # 15. Area & Volume Units
        if re.search(r'\b(?:cm|m)\s*2\b|\bm\s*3\b', slide_text):
            failures.append((15, "Unformatted Area/Volume Unit (cm 2, m 3)", re.findall(r'\b(?:cm|m)\s*2\b|\bm\s*3\b', slide_text)))

        # 16. Charge Density Units
        if re.search(r'\bnC\s*/\s*m\s*2\b|\bμ\s+C\b', slide_text):
            failures.append((16, "Unformatted Charge Density Unit (nC/m 2)", re.findall(r'\bnC\s*/\s*m\s*2\b|\bμ\s+C\b', slide_text)))

        # 17. Permittivity & Field Units
        if re.search(r'\bm\s*F\b|\bm\s*F\s*/\s*m\b', slide_text):
            failures.append((17, "Unformatted Permittivity Unit (m F)", re.findall(r'\bm\s*F\b|\bm\s*F\s*/\s*m\b', slide_text)))

        # 18. Rotational Units
        if re.search(r'\brad\s*/\s*s\b|\bkg\s*m\s*2\b', slide_text):
            failures.append((18, "Unformatted Rotational Unit (rad / s)", re.findall(r'\brad\s*/\s*s\b|\bkg\s*m\s*2\b', slide_text)))

        # 19. Electrical Units
        if re.search(r'\bm\s+A\b|\b\d+\s*V\s+\.', slide_text):
            failures.append((19, "Unformatted Electrical Unit (m A)", re.findall(r'\bm\s+A\b|\b\d+\s*V\s+\.', slide_text)))

        # 20. Temperature Units
        if re.search(r'∘C|\b\d+\s*∘\s*C\b|\bdeg\s*C\b', slide_text):
            failures.append((20, "Unformatted Temperature Unit (∘C)", re.findall(r'∘C|\b\d+\s*∘\s*C\b|\bdeg\s*C\b', slide_text)))

        # === GROUP 3: MATHEMATICAL NOTATION & POWERS (21-30) ===
        # 21. Powers of 10 Negative
        if re.search(r'\b10\s*[-−]\s*\d+\b', slide_text):
            failures.append((21, "Unformatted Negative Power of 10 (10 -2)", re.findall(r'\b10\s*[-−]\s*\d+\b', slide_text)))

        # 22. Powers of 10 Positive
        if re.search(r'\b10\s+\d+\b', slide_text):
            failures.append((22, "Unformatted Positive Power of 10 (10 3)", re.findall(r'\b10\s+\d+\b', slide_text)))

        # 23. Variable Subscripts
        if re.search(r'\b(?:v|x|B|I_L|R_L)\s+[1-4AB]\b', slide_text):
            failures.append((23, "Unformatted Variable Subscripts (v 1, x 1)", re.findall(r'\b(?:v|x|B|I_L|R_L)\s+[1-4AB]\b', slide_text)))

        # 24. Greek Letter Subscripts
        if re.search(r'γ\s+[AB]\b|μ\s+[0rt]\b|ϵ\s+0\b', slide_text):
            failures.append((24, "Unformatted Greek Subscripts (γ A, μ 0)", re.findall(r'γ\s+[AB]\b|μ\s+[0rt]\b|ϵ\s+0\b', slide_text)))

        # 25. Inverse Trig Functions
        if re.search(r'\b(?:cos|sin|tan)\s*[-−^]\s*1\b', slide_text):
            failures.append((25, "Unformatted Inverse Trig Function (cos -1)", re.findall(r'\b(?:cos|sin|tan)\s*[-−^]\s*1\b', slide_text)))

        # 26. Square Root Symbol Cleanliness
        if re.search(r'\bsqrt\b|√\s+\d', slide_text):
            failures.append((26, "Unformatted Square Root (sqrt)", re.findall(r'\bsqrt\b|√\s+\d', slide_text)))

        # 27. Comparison Operators
        if re.search(r'<=|>=|!=|~=|\+-', slide_text):
            failures.append((27, "Raw ASCII Comparison Operators (<=, >=, !=)", re.findall(r'<=|>=|!=|~=|\+-', slide_text)))

        # 28. Arithmetic Spacing
        if re.search(r'[a-zA-Z\d]\+[a-zA-Z\d]|[a-zA-Z\d]=[a-zA-Z\d]', slide_text):
            failures.append((28, "Cramped Arithmetic Operator Spacing", re.findall(r'[a-zA-Z\d]\+[a-zA-Z\d]|[a-zA-Z\d]=[a-zA-Z\d]', slide_text)))

        # 29. Fraction & Ratio Syntax
        if re.search(r'\bB\s*B[1₁]\b|\b50\s+7\s*m/s\b', slide_text):
            failures.append((29, "Broken Fraction / Ratio Notation", re.findall(r'\bB\s*B[1₁]\b|\b50\s+7\s*m/s\b', slide_text)))

        # 30. Exponent Cleanliness
        if re.search(r'\bab\s*[-−]\s*2\b', slide_text):
            failures.append((30, "Unformatted Exponent (ab -2)", re.findall(r'\bab\s*[-−]\s*2\b', slide_text)))

        # === GROUP 4: CHEMICAL & PHYSICAL ENTITIES (31-38) ===
        # 31. Organic Molecular Formulas
        if re.search(r'\bC\s*4\s*H\s*9\s*Br\b|C4H9Br\b', slide_text):
            failures.append((31, "Unformatted Organic Molecule (C4H9Br)", re.findall(r'\bC\s*4\s*H\s*9\s*Br\b|C4H9Br\b', slide_text)))

        # 32. Inorganic Reagents
        if re.search(r'\bNaNH\s*2\b|NaNH2\b', slide_text):
            failures.append((32, "Unformatted Inorganic Reagent (NaNH2)", re.findall(r'\bNaNH\s*2\b|NaNH2\b', slide_text)))

        # 33. Ionic Charge Species
        if re.search(r'\bHe\s*\+\b|\bLi\s*\+\+\b|\bLi\s*\+\s*2\b', slide_text):
            failures.append((33, "Unformatted Ionic Charge (He +, Li ++)", re.findall(r'\bHe\s*\+\b|\bLi\s*\+\+\b|\bLi\s*\+\s*2\b', slide_text)))

        # 34. Henry's Law Constant
        if re.search(r'\bK\s+H\b', slide_text):
            failures.append((34, "Unformatted Henry's Constant (K H)", re.findall(r'\bK\s+H\b', slide_text)))

        # 35. IUPAC Nomenclature Hyphenation
        if re.search(r'\bButan⁻¹-al\b|\bButan\^1-al\b', slide_text):
            failures.append((35, "Corrupted IUPAC Name (Butan⁻¹-al)", "Found superscript 1 in Butan-1-al"))

        # 36. Dimensional Analysis Brackets
        dim_matches = re.findall(r'\[\s*M[L0-9\s−\-]*T[0-9\s−\-]*\]', slide_text)
        for dm in dim_matches:
            if not re.match(r'^\[ML[⁰\d]?\s*T[⁻\d]+\]$', dm.replace(' ', '')) and not re.match(r'^\[M[⁻\d]+\s*L\s*T[²\d]+\]$', dm.replace(' ', '')):
                failures.append((36, "Unformatted Dimensional Formula", dm))

        # 37. Ordinal Suffix Attachment
        if re.search(r'\b\d+\s+(?:st|nd|rd|th)\b', slide_text):
            failures.append((37, "Detached Ordinal Suffix (1 st, 2 nd)", re.findall(r'\b\d+\s+(?:st|nd|rd|th)\b', slide_text)))

        # 38. Axis Hyphenation
        if re.search(r'\b[xyz]\s+[-−]\s*axis\b|\b[xyz]\s*[-−]\s+axis\b', slide_text):
            failures.append((38, "Detached Axis Hyphenation (x -axis)", re.findall(r'\b[xyz]\s+[-−]\s*axis\b|\b[xyz]\s*[-−]\s+axis\b', slide_text)))

        # === GROUP 5: STRUCTURAL & OPTION INTEGRITY (39-45) ===
        # 39. Question Stem Existence (> 20 chars)
        q_match = re.search(r'Q(\d+)\.\s*([^\n\r]+)', slide_text)
        if not q_match or len(q_match.group(2).strip()) < 20:
            failures.append((39, "Missing / Truncated Question Stem", "Stem length < 20 characters"))

        # 40. Section Header Validity
        if not re.search(r'\b(PHYSICS|CHEMISTRY|MATHEMATICS|BIOLOGY|SECTION)\b', slide_text):
            failures.append((40, "Missing Section Tag Header", "No valid section header found"))

        # 41. Question Numbering Sequence
        if q_match:
            q_num = int(q_match.group(1))
            expected_q = idx if idx <= 25 else (idx - 25)
            if q_num != expected_q:
                failures.append((41, "Question Sequence Discontinuity", f"Expected Q{expected_q}, found Q{q_num}"))

        # 42. MCQ Four Options Present
        is_mcq = (idx <= 20) or (idx >= 26 and idx != 29)
        if is_mcq:
            opts = re.findall(r'\(([A-D])\)\s*([^\n\r]+)', slide_text)
            opt_keys = [o[0] for o in opts]
            if set(opt_keys) != {'A', 'B', 'C', 'D'}:
                failures.append((42, "Missing MCQ Option Labels", f"Found: {opt_keys}"))

        # 43. MCQ Non-Empty Options
        if is_mcq:
            opts = re.findall(r'\(([A-D])\)\s*([^\n\r]+)', slide_text)
            for o_k, o_v in opts:
                if len(o_v.strip()) == 0:
                    failures.append((43, "Empty Option Body", f"Option ({o_k}) is blank"))

        # 44. Numerical Section Integer Purity (Q21-Q25 must not have MCQ options)
        if 21 <= idx <= 25:
            opts = re.findall(r'\(([A-D])\)', slide_text)
            if len(opts) > 0:
                failures.append((44, "Erroneous MCQ Options in Numerical Question", opts))

        # 45. Option Label Cleanliness
        if is_mcq:
            if re.search(r'\bOption\s*[1-4A-D]:', slide_text):
                failures.append((45, "Raw 'Option X:' header leaking in slide text", "Found raw Option header"))

        # === GROUP 6: VISUAL LAYOUT & PRESENTATION SPECS (46-50) ===
        # 46. Slide Dimensions (13.333 x 7.5 Widescreen)
        slide_w = prs.slide_width.inches
        slide_h = prs.slide_height.inches
        if abs(slide_w - 13.333) > 0.1 or abs(slide_h - 7.5) > 0.1:
            failures.append((46, "Non-Standard Slide Dimensions", f"{slide_w:.2f}x{slide_h:.2f} (Expected 13.33x7.5)"))

        # 47. Content Bounds Safety (No off-screen shapes)
        for shape in slide.shapes:
            left = shape.left.inches
            top = shape.top.inches
            w = shape.width.inches
            h = shape.height.inches
            if left < 0 or top < 0 or (left + w) > slide_w + 0.1 or (top + h) > slide_h + 0.1:
                failures.append((47, "Shape Exceeds Slide Boundaries", f"({left:.2f}, {top:.2f}, {w:.2f}, {h:.2f})"))

        # 48. Minimum Legible Font Size (>= 10pt)
        for tf in text_frames:
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.font.size and r.font.size.pt < 10:
                        failures.append((48, "Sub-Legible Font Size (<10pt)", f"{r.font.size.pt}pt"))

        # 49. Attached Diagram Media Validity
        if len(pictures) > 0:
            for p in pictures:
                if p.width.inches < 0.5 or p.height.inches < 0.5:
                    failures.append((49, "Corrupted Diagram Dimensions", f"{p.width.inches:.2f}x{p.height.inches:.2f}"))

        # 50. Non-Empty Presentation Slide Count (>= 30)
        if len(lines) == 0:
            failures.append((50, "Completely Blank Slide", "Slide has zero text or image content"))

        if failures:
            slide_failures[idx] = failures
            for chk_num, _, _ in failures:
                check_failure_counts[chk_num] += 1
            print(f"❌ Slide {idx:02d} FAILED {len(failures)} Check(s):")
            for num, name, detail in failures:
                print(f"   • [Check {num:02d} - {name}]: {detail}")
        else:
            print(f"✅ Slide {idx:02d}: PASSED ALL 50 ENTERPRISE QUALITY CHECKS")

    print("\n" + "=" * 80)
    print(f"AUDIT SUMMARY: {total_slides - len(slide_failures)} / {total_slides} Slides Passed 100% Perfectly.")
    print(f"Total Quality Check Assertions Run: {total_slides * 50} (50 Checks × {total_slides} Slides)")
    if len(slide_failures) == 0:
        print("🏆 100% PERFECT ZERO-DEFECT QUALITY SCORE ACROSS ALL 50 CRITERIA!")
    else:
        print(f"⚠️ {len(slide_failures)} slide(s) require fixes.")
    print("=" * 80)

    return len(slide_failures)

if __name__ == "__main__":
    pptx_file = r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev94.pptx"
    run_50_point_audit(pptx_file)

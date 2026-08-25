import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev87.pptx")
print("=== SLIDE 6 (INDEX 5) ===")
print(prs.slides[5].shapes[0].text_frame.text)

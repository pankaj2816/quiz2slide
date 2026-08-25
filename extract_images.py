import pptx
import os

prs = pptx.Presentation(r'd:\Work\kaku\test_output\Quiz_Presentation_30_Slides.pptx')
os.makedirs(r'd:\Work\kaku\pdf_quiz_gh_pages\extracted_slide_images', exist_ok=True)

img_count = 0
for idx, s in enumerate(prs.slides, 1):
    for shp_idx, sp in enumerate(s.shapes):
        if sp.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            img_count += 1
            image = sp.image
            image_bytes = image.blob
            img_filename = f"slide_{idx}_img_{shp_idx}.{image.ext}"
            img_path = os.path.join(r'd:\Work\kaku\pdf_quiz_gh_pages\extracted_slide_images', img_filename)
            with open(img_path, 'wb') as f:
                f.write(image_bytes)
            print(f"Slide {idx}: saved {img_filename} (size: {len(image_bytes)} bytes)")

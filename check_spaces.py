import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
import re

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev87.pptx")

print("=== CHECKING MULTIPLE SPACES ACROSS ALL SLIDES ===")
for idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text
                double_spaces = re.findall(r' {2,}', txt)
                if double_spaces:
                    print(f"Slide {idx:02d} Paragraph has multiple spaces:")
                    print(f"  RAW: {repr(txt)}")

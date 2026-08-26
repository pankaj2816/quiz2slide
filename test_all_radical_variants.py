import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.5))
tf = txBox.text_frame
tf.word_wrap = True

variants = [
    ("A. Combining Overline (U+0305)", "cos⁻¹(1/√3\u0305)", "cos⁻¹(√2\u0305/3)"),
    ("B. Combining Macron (U+0304)", "cos⁻¹(1/√3\u0304)", "cos⁻¹(√2\u0304/3)"),
    ("C. Overline Character (U+203E)", "cos⁻¹(1/√‾3)", "cos⁻¹(√‾2/3)"),
    ("D. Macron Character (U+00AF)", "cos⁻¹(1/√¯3)", "cos⁻¹(√¯2/3)"),
    ("E. Full Radicand Parentheses", "cos⁻¹(1/√(3))", "cos⁻¹(√(2)/3)"),
    ("F. Full Radicand Bracket", "cos⁻¹(1/√[3])", "cos⁻¹(√[2]/3)"),
    ("G. Spaced Vinculum Overline", "cos⁻¹(1/√ 3\u0305)", "cos⁻¹(√ 2\u0305/3)"),
]

for name, v1, v2 in variants:
    p = tf.add_paragraph()
    p.text = f"{name}:  (A) {v1}    (D) {v2}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

prs.save(r"d:\Work\kaku\test_output\all_radical_variants.pptx")
print("Saved all_radical_variants.pptx")

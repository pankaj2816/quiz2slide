import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev93.pptx")
print(f"=== FULL PRESENTATION REPORT: {len(prs.slides)} SLIDES ===")

for idx, slide in enumerate(prs.slides, 1):
    texts = []
    has_image = any(shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    texts.append(t)
    full_content = " | ".join(texts)
    img_tag = "[DIAGRAM]" if has_image else "[TEXT-ONLY]"
    print(f"Slide {idx:02d} {img_tag}: {full_content[:110]}...")

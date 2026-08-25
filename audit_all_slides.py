import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
import re

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Verified.pptx")

print("=== SEARCHING FOR UNFORMATTED UNITS & POWERS ACROSS ALL 30 SLIDES ===")

for idx, slide in enumerate(prs.slides, 1):
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            slide_text += "\n" + shape.text_frame.text
            
    # Check for raw exponents, loose signs, spacing in units
    patterns = [
        r'\b\w+\s*[-−]\s*\d+\b', # e.g. h -1, ms -2, ab -2
        r'\b10\s*[-−]\s*\d+\b',  # e.g. 10 -2
        r'\b10\s+\d+\b',         # e.g. 10 3, 10 5, 10 11
        r'\b(?:cm|kg|m|s)\s*\d+\b', # e.g. cm 2, kg 3
        r'\b(?:km|m)\s+h\s*[-−]?\s*1\b', # e.g. km h -1
        r'\b[A-Za-z]\s+[1-4]\b', # e.g. v 1, x 1, B 1
        r'\b[A-Za-z]\s+[\+\-]\b', # e.g. He +, Li ++
        r'\b\d+\s+(?:st|nd|rd|th)\b', # e.g. 1 st, 2 nd
        r'\b\d+\s*°C?\b',
        r'/\s*m\b',
        r'F/\s*m\b',
        r'B\s*B₁',
        r'γ\s*A',
        r'K\s*H'
    ]
    
    found = []
    for pat in patterns:
        matches = re.findall(pat, slide_text)
        if matches:
            found.extend(matches)
            
    print(f"Slide {idx:02d}: Found potential issues -> {found}")
    print("Full Slide Text:")
    print(slide_text.strip())
    print("=" * 60)

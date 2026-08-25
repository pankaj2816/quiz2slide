import sys
sys.stdout.reconfigure(encoding='utf-8')
import os
import time
import base64
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import pptx

pdf_path = r"C:\Users\Lenovo\.gemini\antigravity\brain\8ad26341-31c8-489f-8b94-35a6e7a0c5aa\.user_uploaded\media_1787675766400.pdf"
html_path = r"d:\Work\kaku\pdf_quiz_gh_pages\index.html"
output_pptx = r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev87.pptx"

chrome_options = Options()
chrome_options.add_argument("--headless=new")
chrome_options.add_argument("--disable-gpu")
chrome_options.add_argument("--window-size=1920,1080")
chrome_options.add_argument("--allow-file-access-from-files")
chrome_options.add_argument("--disable-web-security")

driver = webdriver.Chrome(options=chrome_options)

try:
    print(f"Loading local web app: file:///{html_path}")
    driver.get(f"file:///{html_path}")

    # Wait for page to load
    WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.ID, "pdfFileInput")))
    rev_badge = driver.find_element(By.CLASS_NAME, "revision-badge").text
    print(f"Loaded Page. Revision badge: {rev_badge}")

    # Upload PDF
    print(f"Uploading PDF: {pdf_path}")
    file_input = driver.find_element(By.ID, "pdfFileInput")
    file_input.send_keys(pdf_path)

    # Click Parse PDF
    time.sleep(1)
    parse_btn = driver.find_element(By.ID, "btnParsePdf")
    parse_btn.click()
    print("Clicked 'Parse PDF Questions'. Waiting for parsing to complete...")

    # Wait for parsing to finish
    WebDriverWait(driver, 30).until(lambda d: len(d.find_elements(By.CLASS_NAME, "question-item")) > 0)
    
    q_items = driver.find_elements(By.CLASS_NAME, "question-item")
    print(f"Extracted {len(q_items)} questions!")

    # Check which questions have diagrams attached in the UI
    for idx, item in enumerate(q_items, 1):
        has_diagram = len(item.find_elements(By.CLASS_NAME, "badge-diagram")) > 0
        stem = item.find_element(By.CLASS_NAME, "q-stem").text.strip().replace('\n', ' ')
        if has_diagram:
            print(f"  [DIAGRAM] Question {idx}: {stem[:80]}...")

    # Generate PowerPoint via JS execution to get the exact blob base64
    print("\nTriggering PowerPoint generation via PptxGenJS...")
    driver.execute_script("""
        window._generatedPptxBase64 = null;
        const btn = document.getElementById('btnGeneratePpt');
        // Override showSaveFilePicker so it doesn't prompt in headless mode
        window.showSaveFilePicker = undefined;
        // Intercept downloadBlobDirectly
        window.downloadBlobDirectly = function(blob, filename) {
            const reader = new FileReader();
            reader.onloadend = function() {
                window._generatedPptxBase64 = reader.result;
            };
            reader.readAsDataURL(blob);
        };
        btn.click();
    """)

    # Wait for generated PPTX Base64
    b64_data = None
    for _ in range(30):
        time.sleep(1)
        b64_data = driver.execute_script("return window._generatedPptxBase64;")
        if b64_data:
            break

    if not b64_data:
        raise Exception("Timed out waiting for PPTX generation in browser.")

    # Save to disk
    header, b64_str = b64_data.split(',', 1)
    pptx_bytes = base64.b64decode(b64_str)
    with open(output_pptx, 'wb') as f:
        f.write(pptx_bytes)

    print(f"Successfully generated and saved: {output_pptx} ({len(pptx_bytes)} bytes)\n")

finally:
    driver.quit()

# Inspect generated PPTX in detail
print("==================== ANALYZING FINAL PPTX OUTPUT ====================")
prs = pptx.Presentation(output_pptx)
print(f"Total Slides: {len(prs.slides)}")

for idx, s in enumerate(prs.slides, 1):
    images = [sp for sp in s.shapes if sp.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE]
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    print(f"\n--- SLIDE {idx} (Images: {len(images)}) ---")
    for t in texts:
        print(t)

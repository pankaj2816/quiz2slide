import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r'd:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Verified.pptx')

for i in range(1, 6):
    s = prs.slides[i - 1]
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    print(f"\n==================== SLIDE {i} ====================")
    for t in texts:
        print(t)

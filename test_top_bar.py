import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(6.5))
tf = txBox.text_frame
tf.word_wrap = True

# Add title
p = tf.paragraphs[0]
p.text = "SQUARE ROOT TOP BAR COMPARISON IN POWERPOINT"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 215, 0)

tests = [
    ("1. Overline Symbol \u203E (High Top Bar):", "cos⁻¹(1/√‾3)", "cos⁻¹(√‾2/3)"),
    ("2. Macron Symbol \u00AF (Medium Top Bar):", "cos⁻¹(1/√¯3)", "cos⁻¹(√¯2/3)"),
    ("3. Overline with Paren \u203E(3):", "cos⁻¹(1/√‾(3))", "cos⁻¹(√‾(2)/3)"),
    ("4. Combining Overline \u0305 (Glued Top Bar):", "cos⁻¹(1/√3\u0305)", "cos⁻¹(√2\u0305/3)"),
    ("5. Root with Overline Character:", "√‾3 m", "10√‾2 N")
]

for title, optA, optD in tests:
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    
    p2 = tf.add_paragraph()
    p2.text = f"   (A) {optA}       (D) {optD}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(255, 255, 255)

# Also add OMML Native Math Equation Box
omml_sp = parse_xml("""
<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
      xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
      xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">
  <p:nvSpPr>
    <p:cNvPr id="200" name="Native Math Equation"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="1000000" y="5500000"/>
      <a:ext cx="9000000" cy="1500000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square"/>
    <a:lstStyle/>
    <a:p>
      <m:oMathPara>
        <m:oMath>
          <m:r><m:rPr><m:scr m:val="roman"/></m:rPr><m:t>6. Native OMML 2D Vector Radical:   (A) cos</m:t></m:r>
          <m:sSup><m:e><m:r><m:t></m:t></m:r></m:e><m:sup><m:r><m:t>-1</m:t></m:r></m:sup></m:sSup>
          <m:r><m:t>(1/</m:t></m:r>
          <m:rad><m:radPr><m:degHide m:val="1"/></m:radPr><m:deg/><m:e><m:r><m:t>3</m:t></m:r></m:e></m:rad>
          <m:r><m:t>)        (D) cos</m:t></m:r>
          <m:sSup><m:e><m:r><m:t></m:t></m:r></m:e><m:sup><m:r><m:t>-1</m:t></m:r></m:sup></m:sSup>
          <m:r><m:t>(</m:t></m:r>
          <m:rad><m:radPr><m:degHide m:val="1"/></m:radPr><m:deg/><m:e><m:r><m:t>2</m:t></m:r></m:e></m:rad>
          <m:r><m:t>/3)</m:t></m:r>
        </m:oMath>
      </m:oMathPara>
    </a:p>
  </p:txBody>
</p:sp>
""")

slide.shapes._spTree.append(omml_sp)
prs.save(r"d:\Work\kaku\test_output\top_bar_comparison.pptx")
print("Saved top_bar_comparison.pptx successfully!")

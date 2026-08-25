import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
import os

prs = pptx.Presentation(r'd:\Work\kaku\test_output\Quiz_Presentation_6_Slides.pptx')
print(f"Total slides: {len(prs.slides)}")

for idx, s in enumerate(prs.slides, 1):
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    pics = [sp for sp in s.shapes if sp.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE]
    print(f"\n==================== SLIDE {idx} (Pictures: {len(pics)}) ====================")
    for t in texts:
        print(t)
        print("-" * 40)

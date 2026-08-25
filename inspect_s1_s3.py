import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Verified.pptx")

for idx in range(3):
    slide = prs.slides[idx]
    print(f"\n=== SLIDE {idx+1} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text.strip())

import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add text box
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Testing Native PowerPoint OMML Equation Radical:"
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(255, 255, 255)

# OMML XML for cos^-1(1/\sqrt{3})
omml_xml = """
<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
      xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
      xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">
  <p:nvSpPr>
    <p:cNvPr id="100" name="Equation 1"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="1000000" y="2500000"/>
      <a:ext cx="8000000" cy="2000000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" rtlCol="0"/>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="l"/>
      <m:oMathPara>
        <m:oMath>
          <m:r><m:rPr><m:scr m:val="roman"/><m:sty m:val="p"/></m:rPr><m:t>(A) cos</m:t></m:r>
          <m:sSup>
            <m:sSupPr><m:ctrlPr/></m:sSupPr>
            <m:e><m:r><m:t></m:t></m:r></m:e>
            <m:sup><m:r><m:t>-1</m:t></m:r></m:sup>
          </m:sSup>
          <m:r><m:t>(1/</m:t></m:r>
          <m:rad>
            <m:radPr><m:degHide m:val="1"/></m:radPr>
            <m:deg/>
            <m:e><m:r><m:t>3</m:t></m:r></m:e>
          </m:rad>
          <m:r><m:t>)</m:t></m:r>
        </m:oMath>
      </m:oMathPara>
    </a:p>
  </p:txBody>
</p:sp>
"""

try:
    sp = parse_xml(omml_xml)
    slide.shapes._spTree.append(sp)
    prs.save(r"d:\Work\kaku\test_output\omml_math_test.pptx")
    print("✅ Successfully generated omml_math_test.pptx with native 2D Vector Radical!")
except Exception as e:
    print(f"❌ Error: {e}")

import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx
import re
import os

def run_25_point_audit(pptx_path):
    print("=" * 75)
    print("         EXHAUSTIVE 25-POINT QUALITY ASSURANCE & INTEGRITY AUDIT")
    print(f"         Target File: {pptx_path}")
    print("=" * 75)

    if not os.path.exists(pptx_path):
        print(f"❌ ERROR: PPTX file not found at {pptx_path}")
        return 1

    prs = pptx.Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"Total Slides Loaded: {total_slides}\n")

    slide_failures = {}

    # Define the 25 Check Functions
    def check_slide(idx, slide):
        failures = []
        slide_text = ""
        text_frames = []
        pictures = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frames.append(shape.text_frame)
                slide_text += "\n" + shape.text_frame.text
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                pictures.append(shape)

        lines = [l for l in slide_text.split('\n') if l.strip()]

        # 1. Consecutive Whitespace
        if re.search(r' {2,}', slide_text):
            failures.append(("1. Consecutive Whitespace", re.findall(r' {2,}', slide_text)))

        # 2. Stray Space Before Punctuation
        if re.search(r' [.,:;?!%°]', slide_text):
            failures.append(("2. Stray Space Before Punctuation", re.findall(r' [.,:;?!%°]', slide_text)))

        # 3. Missing Space After Punctuation (ignoring decimals like 1.5, 0.2)
        missing_space = re.findall(r'[,;?!](?=[A-Za-z])', slide_text)
        if missing_space:
            failures.append(("3. Missing Space After Punctuation", missing_space))

        # 4. Stray Padding Inside Parentheses
        if re.search(r'\(\s+[^\s]|\s+\)', slide_text):
            failures.append(("4. Stray Space in Parentheses", re.findall(r'\(\s+[^\s]|\s+\)', slide_text)))

        # 5. Scientific Unit Negative Exponents
        if re.search(r'\b(?:km\s+h|ms|Nm|m\s+s)\s*[-−]\s*\d+\b', slide_text):
            failures.append(("5. Unformatted Unit Negative Exponents", re.findall(r'\b(?:km\s+h|ms|Nm|m\s+s)\s*[-−]\s*\d+\b', slide_text)))

        # 6. Unformatted Powers of 10
        if re.search(r'\b10\s*[-−]\s*\d+\b|\b10\s+\d+\b', slide_text):
            failures.append(("6. Unformatted Powers of 10", re.findall(r'\b10\s*[-−]\s*\d+\b|\b10\s+\d+\b', slide_text)))

        # 7. Unformatted Variable Subscripts
        if re.search(r'\b(?:v|x|B|I_L|R_L|K_H)\s+[1-4AB]\b', slide_text):
            failures.append(("7. Unformatted Variable Subscripts", re.findall(r'\b(?:v|x|B|I_L|R_L|K_H)\s+[1-4AB]\b', slide_text)))

        # 8. Greek Letter Subscripts & Symbols
        if re.search(r'γ\s+[AB]\b|μ\s+[0rt]\b|ϵ\s+0\b', slide_text):
            failures.append(("8. Unformatted Greek Subscripts", re.findall(r'γ\s+[AB]\b|μ\s+[0rt]\b|ϵ\s+0\b', slide_text)))

        # 9. Chemical Formula Notation
        if re.search(r'\b(?:C4H9Br|NaNH2|He\s*\+|Li\s*\+\+)\b', slide_text):
            failures.append(("9. Unformatted Chemical Formulas", re.findall(r'\b(?:C4H9Br|NaNH2|He\s*\+|Li\s*\+\+)\b', slide_text)))

        # 10. Dimensional Formula Notation
        dim_matches = re.findall(r'\[\s*M[L0-9\s−\-]*T[0-9\s−\-]*\]', slide_text)
        for dm in dim_matches:
            if not re.match(r'^\[ML[⁰\d]?\s*T[⁻\d]+\]$', dm.replace(' ', '')) and not re.match(r'^\[M[⁻\d]+\s*L\s*T[²\d]+\]$', dm.replace(' ', '')):
                failures.append(("10. Unformatted Dimensional Formula", dm))

        # 11. Ordinal Numbers
        if re.search(r'\b\d+\s+(?:st|nd|rd|th)\b', slide_text):
            failures.append(("11. Unformatted Ordinals", re.findall(r'\b\d+\s+(?:st|nd|rd|th)\b', slide_text)))

        # 12. Raw Degree Signs
        if re.search(r'∘C|\b\d+\s*∘\b', slide_text):
            failures.append(("12. Raw Degree Signs", re.findall(r'∘C|\b\d+\s*∘\b', slide_text)))

        # 13. Inverse Trig Notation
        if re.search(r'\b(?:cos|sin|tan)\s*[-−^]\s*1\b', slide_text):
            failures.append(("13. Unformatted Inverse Trig", re.findall(r'\b(?:cos|sin|tan)\s*[-−^]\s*1\b', slide_text)))

        # 14. Comparison Operators
        if re.search(r'<=|>=|!=|~=|\+-', slide_text):
            failures.append(("14. Raw ASCII Math Operators", re.findall(r'<=|>=|!=|~=|\+-', slide_text)))

        # 15. Reaction / Logic Arrows
        if re.search(r'-->|->|<--|<-|<->', slide_text):
            failures.append(("15. Raw ASCII Arrows", re.findall(r'-->|->|<--|<-|<->', slide_text)))

        # 16. Broken Ratio Expressions
        if re.search(r'\bB\s*B[1₁]\b|\b50\s+7\s*m/s\b', slide_text):
            failures.append(("16. Broken Ratio Expressions", re.findall(r'\bB\s*B[1₁]\b|\b50\s+7\s*m/s\b', slide_text)))

        # 17. Loose Negative Exponents
        if re.search(r'\bab\s*[-−]\s*2\b', slide_text):
            failures.append(("17. Loose Negative Exponents", re.findall(r'\bab\s*[-−]\s*2\b', slide_text)))

        # 18. Question Stem Integrity (Length > 15 chars)
        q_match = re.search(r'Q(\d+)\.\s*([^\n\r]+)', slide_text)
        if not q_match or len(q_match.group(2).strip()) < 15:
            failures.append(("18. Question Stem Integrity", "Question stem is missing or too short"))

        # 19. Section Header Existence
        if not re.search(r'\b(PHYSICS|CHEMISTRY|MATHEMATICS|BIOLOGY|SECTION)\b', slide_text):
            failures.append(("19. Missing Section Header", "No subject section found on slide"))

        # 20. Question Numbering Sequence
        if q_match:
            q_num = int(q_match.group(1))
            expected_q = idx if idx <= 25 else (idx - 25) # Physics 1-25, Chemistry 1-5
            if q_num != expected_q:
                failures.append(("20. Question Number Sequence Mismatch", f"Expected Q{expected_q}, got Q{q_num}"))

        # 21. Multiple Choice Options Completeness (For Q1-Q20 and Q26-Q30)
        is_mcq = (idx <= 20) or (idx >= 26 and idx != 29) # Q29 is graph matching
        if is_mcq:
            opts = re.findall(r'\(([A-D])\)\s*([^\n\r]+)', slide_text)
            opt_keys = [o[0] for o in opts]
            if set(opt_keys) != {'A', 'B', 'C', 'D'}:
                failures.append(("21. Incomplete Options Set (A, B, C, D)", f"Found options: {opt_keys}"))
            for o_key, o_val in opts:
                if len(o_val.strip()) == 0:
                    failures.append(("21. Empty Option Content", f"Option ({o_key}) is empty"))

        # 22. Numerical Section Integer Purity (Q21-Q25 should NOT have A, B, C, D)
        if 21 <= idx <= 25:
            opts = re.findall(r'\(([A-D])\)', slide_text)
            if len(opts) > 0:
                failures.append(("22. Erroneous Options in Numerical Question", opts))

        # 23. Layout Bounds Check (Within 13.333 x 7.5 in)
        slide_w = prs.slide_width.inches
        slide_h = prs.slide_height.inches
        for shape in slide.shapes:
            left = shape.left.inches
            top = shape.top.inches
            w = shape.width.inches
            h = shape.height.inches
            if left < 0 or top < 0 or (left + w) > slide_w + 0.1 or (top + h) > slide_h + 0.1:
                failures.append(("23. Layout Out of Bounds", f"Shape at ({left:.2f}, {top:.2f}, {w:.2f}, {h:.2f}) exceeds slide ({slide_w:.2f}x{slide_h:.2f})"))

        # 24. Diagram Image Validity
        if len(pictures) > 0:
            for p in pictures:
                if p.width.inches < 0.5 or p.height.inches < 0.5:
                    failures.append(("24. Corrupted/Zero-Dimension Image", f"Image size {p.width.inches:.2f}x{p.height.inches:.2f}"))

        # 25. Typography Hierarchy & Text Readability
        for tf in text_frames:
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.font.size and r.font.size.pt < 9:
                        failures.append(("25. Illegible Font Size (<9pt)", f"Font size {r.font.size.pt}pt in '{r.text[:20]}'"))

        return failures

    for idx, slide in enumerate(prs.slides, 1):
        failures = check_slide(idx, slide)
        if failures:
            slide_failures[idx] = failures
            print(f"❌ Slide {idx:02d} FAILED {len(failures)} Check(s):")
            for name, detail in failures:
                print(f"   • [{name}]: {detail}")
        else:
            print(f"✅ Slide {idx:02d}: PASSED ALL 25 QUALITY CHECKS")

    print("\n" + "=" * 75)
    print(f"FINAL AUDIT RESULT: {total_slides - len(slide_failures)} / {total_slides} Slides Passed 100% Perfectly.")
    if len(slide_failures) == 0:
        print("🏆 100% PERFECT ZERO-DEFECT QUALITY SCORE ACROSS ALL 25 CRITERIA!")
    else:
        print(f"⚠️ {len(slide_failures)} slide(s) need adjustments.")
    print("=" * 75)

    return len(slide_failures)

if __name__ == "__main__":
    pptx_path = r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev89.pptx"
    run_25_point_audit(pptx_path)

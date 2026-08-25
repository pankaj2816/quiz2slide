import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r'd:\Work\kaku\test_output\Quiz_Presentation_30_Slides.pptx')
print(f"Total slides: {len(prs.slides)}")

for idx, s in enumerate(prs.slides, 1):
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    print(f"\n==================== SLIDE {idx} ====================")
    for t in texts:
        print(t)

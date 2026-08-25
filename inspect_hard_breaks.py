import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Rev89.pptx")

print("=== INSPECTING STEM HARD BREAKS ACROSS ALL SLIDES ===")

for idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            p_texts = [p.text for p in shape.text_frame.paragraphs]
            print(f"\n--- SLIDE {idx:02d} Paragraph Count: {len(p_texts)} ---")
            for p_i, p_t in enumerate(p_texts):
                print(f"  P{p_i}: {repr(p_t)}")

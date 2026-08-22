import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r'd:\Work\kaku\test_output\Quiz_Presentation_96_Slides.pptx')

for i in range(1, 10):
    s = prs.slides[i - 1]
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    print(f"Slide {i}: {texts[0][:120] if texts else 'EMPTY'}...")

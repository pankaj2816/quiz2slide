import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r"d:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Verified.pptx")

for idx in range(17):
    slide = prs.slides[idx]
    print(f"\n========================================================")
    print(f"                     SLIDE {idx+1:02d}")
    print(f"========================================================")
    
    images_count = sum(1 for shape in slide.shapes if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE)
    print(f"Attached Images/Diagrams: {images_count}")
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text.strip())
            print("-" * 40)

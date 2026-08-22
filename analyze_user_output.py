import sys
sys.stdout.reconfigure(encoding='utf-8')
import os
import pptx
from pptx.util import Inches, Pt

pptx_path = r'd:\Work\kaku\test_output\Quiz_Presentation_96_Slides.pptx'
prs = pptx.Presentation(pptx_path)

print(f"Total slides in presentation: {len(prs.slides)}")
print(f"Slide dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")

# Analyze each slide
issues = []
for idx, slide in enumerate(prs.slides, 1):
    text_content = []
    has_image = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_content.append(shape.text_frame.text)
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            has_image = True
            
    full_slide_text = "\n".join(text_content)
    
    # Check checks
    if f"Q{idx}." not in full_slide_text and f"{idx}." not in full_slide_text:
        issues.append(f"Slide {idx}: Missing question number header")
    if "(A)" not in full_slide_text or "(B)" not in full_slide_text:
        issues.append(f"Slide {idx}: Missing option A or B")
    if idx == 22 and not has_image:
        issues.append("Slide 22: Expected image/diagram not found as picture shape")
    if idx != 22 and has_image:
        issues.append(f"Slide {idx}: Extraneous picture shape found")

print(f"Automated check found {len(issues)} structural issues.")
for iss in issues[:10]:
    print("  -", iss)

print("\nSampling Slide Texts:")
for s_num in [1, 10, 11, 22, 27, 62, 96]:
    s = prs.slides[s_num - 1]
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    print(f"\n--- SLIDE {s_num} ---")
    print("\n".join(texts)[:250] + "...")

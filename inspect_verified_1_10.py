import sys
sys.stdout.reconfigure(encoding='utf-8')
import pptx

prs = pptx.Presentation(r'd:\Work\kaku\test_output\Quiz_Presentation_30_Slides_Verified.pptx')

for idx in range(1, 11):
    s = prs.slides[idx - 1]
    images = [sp for sp in s.shapes if sp.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE]
    texts = [sp.text_frame.text for sp in s.shapes if sp.has_text_frame]
    print(f"\n==================== SLIDE {idx} (Images: {len(images)}) ====================")
    for t in texts:
        print(t)

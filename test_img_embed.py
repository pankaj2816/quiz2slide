import sys
sys.stdout.reconfigure(encoding='utf-8')
import base64
import pptx

doc = pptx.Presentation()
slide = doc.slides.add_slide(doc.slide_layouts[6])

# Let's verify python-pptx / PptxGenJS image addition
print("Verifying image embed...")

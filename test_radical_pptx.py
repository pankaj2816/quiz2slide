import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Q1. Radical Overline Test:\n\n(A) cos⁻¹(1/√3\u0305)\n(B) cos⁻¹(2/3)\n(C) cos⁻¹(1/3)\n(D) cos⁻¹(√2\u0305/3)\n\nQ6: 112.5 × √3\u0305 m\nQ9: x₁ = √7\u0305 sin 5t cm, x₂ = 2√7\u0305 sin(5t + π/3) cm\nQ13: 10√2\u0305 N"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)

prs.save(r"d:\Work\kaku\test_output\radical_test.pptx")
print("Saved radical_test.pptx successfully!")

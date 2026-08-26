import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(6.0))
tf = txBox.text_frame
tf.word_wrap = True

styles = [
    ("Style 1: Combining Overline (Current)", "cos⁻¹(1/√3\u0305)", "cos⁻¹(√2\u0305/3)"),
    ("Style 2: High Vinculum Overscore (√‾)", "cos⁻¹(1/√‾3)", "cos⁻¹(√‾2/3)"),
    ("Style 3: Spaced Combining Vinculum", "cos⁻¹(1/√\u200a3\u0305)", "cos⁻¹(√\u200a2\u0305/3)"),
    ("Style 4: Parenthesized Clean Radical √(3)", "cos⁻¹(1/√(3))", "cos⁻¹(√(2)/3)"),
    ("Style 5: Bracketed Clean Radical √[3]", "cos⁻¹(1/√[3])", "cos⁻¹(√[2]/3)"),
    ("Style 6: High Macron ¯ (√¯)", "cos⁻¹(1/√¯3)", "cos⁻¹(√¯2/3)")
]

for title, optA, optD in styles:
    p = tf.add_paragraph()
    p.text = f"--- {title} ---"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)
    
    p2 = tf.add_paragraph()
    p2.text = f"  (A) {optA}       (D) {optD}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    
    p_sp = tf.add_paragraph()
    p_sp.text = ""
    p_sp.font.size = Pt(6)

prs.save(r"d:\Work\kaku\test_output\radical_styles_comparison.pptx")
print("Saved radical_styles_comparison.pptx successfully!")

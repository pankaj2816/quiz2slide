import sys
sys.stdout.reconfigure(encoding='utf-8')
import os
import time
import base64
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import pptx

pdf_path = r"C:\Users\Lenovo\.gemini\antigravity\brain\8ad26341-31c8-489f-8b94-35a6e7a0c5aa\.user_uploaded\media_1787675766400.pdf"
html_path = r"d:\Work\kaku\pdf_quiz_gh_pages\index.html"

def test_scaling_preset(scale_val, output_name):
    output_pptx = os.path.join(r"d:\Work\kaku\test_output", output_name)
    chrome_options = Options()
    chrome_options.add_argument("--headless=new")
    chrome_options.add_argument("--disable-gpu")
    chrome_options.add_argument("--window-size=1920,1080")
    chrome_options.add_argument("--allow-file-access-from-files")
    chrome_options.add_argument("--disable-web-security")

    driver = webdriver.Chrome(options=chrome_options)
    try:
        driver.get(f"file:///{html_path}")
        WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.ID, "pdfFileInput")))

        # Upload PDF
        file_input = driver.find_element(By.ID, "pdfFileInput")
        file_input.send_keys(pdf_path)

        # Click Parse PDF
        time.sleep(1)
        parse_btn = driver.find_element(By.ID, "btnParsePdf")
        parse_btn.click()

        WebDriverWait(driver, 30).until(lambda d: len(d.find_elements(By.CLASS_NAME, "question-item")) > 0)

        # Set font scale slider
        driver.execute_script(f"""
            const slider = document.getElementById('inputFontScale');
            slider.value = {scale_val};
            slider.dispatchEvent(new Event('input'));
        """)
        time.sleep(0.5)

        # Hook downloadBlobDirectly
        driver.execute_script("""
            window._generatedPptxBase64 = null;
            window.showSaveFilePicker = undefined;
            window.downloadBlobDirectly = function(blob, filename) {
                const reader = new FileReader();
                reader.onloadend = function() {
                    window._generatedPptxBase64 = reader.result;
                };
                reader.readAsDataURL(blob);
            };
            document.getElementById('btnGeneratePpt').click();
        """)

        b64_data = None
        for _ in range(30):
            time.sleep(1)
            b64_data = driver.execute_script("return window._generatedPptxBase64;")
            if b64_data:
                break

        if not b64_data:
            print(f"❌ Failed to get PPTX blob for scale {scale_val}%", flush=True)
            return False

        header, b64_str = b64_data.split(',', 1)
        pptx_bytes = base64.b64decode(b64_str)
        with open(output_pptx, 'wb') as f:
            f.write(pptx_bytes)

        prs = pptx.Presentation(output_pptx)
        print(f"✅ Scale {scale_val}%: Generated {output_name} ({len(prs.slides)} slides, {len(pptx_bytes)/1024:.1f} KB)", flush=True)

        # Inspect slide 2 font size
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            print(f"   [Scale {scale_val}%] '{r.text[:30]}' -> {r.font.size.pt}pt", flush=True)
        return True
    finally:
        driver.quit()

if __name__ == "__main__":
    print("=== TESTING GLOBAL FONT SCALING (COMPACT 85%, STANDARD 100%, LARGE 120%, JUMBO 140%) ===", flush=True)
    test_scaling_preset(85, "Quiz_Presentation_Compact_85pct.pptx")
    test_scaling_preset(100, "Quiz_Presentation_Standard_100pct.pptx")
    test_scaling_preset(120, "Quiz_Presentation_Large_120pct.pptx")
    test_scaling_preset(140, "Quiz_Presentation_Jumbo_140pct.pptx")
